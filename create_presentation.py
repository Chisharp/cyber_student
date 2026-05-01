"""
Generates GDPR_Screencast.pptx — 6-slide presentation for the screencast.
Run: python create_presentation.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Colour palette ────────────────────────────────────────────────────────────
DARK_BG    = RGBColor(0x1E, 0x1E, 0x2E)   # dark navy
ACCENT     = RGBColor(0x00, 0xAE, 0xEF)   # cyan blue
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xCC, 0xCC, 0xCC)
GREEN      = RGBColor(0x00, 0xC8, 0x5A)
YELLOW     = RGBColor(0xFF, 0xD7, 0x00)
CODE_BG    = RGBColor(0x2D, 0x2D, 0x3F)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def set_bg(slide, colour: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = colour


def add_textbox(slide, text, left, top, width, height,
                font_size=18, bold=False, colour=WHITE,
                align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = colour
    return txBox


def add_bullet_box(slide, items, left, top, width, height,
                   font_size=17, colour=WHITE, bullet_colour=ACCENT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        # bullet symbol
        run_b = p.add_run()
        run_b.text = "▸  "
        run_b.font.size = Pt(font_size)
        run_b.font.color.rgb = bullet_colour
        run_b.font.bold = True
        # text
        run_t = p.add_run()
        run_t.text = item
        run_t.font.size = Pt(font_size)
        run_t.font.color.rgb = colour
        p.space_after = Pt(6)


def add_code_box(slide, code, left, top, width, height, font_size=13):
    """Dark code block."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = CODE_BG
    shape.line.color.rgb = ACCENT
    shape.line.width = Pt(1)

    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = code
    run.font.size = Pt(font_size)
    run.font.color.rgb = GREEN
    run.font.name = "Courier New"


def add_accent_bar(slide, top=Inches(1.15)):
    bar = slide.shapes.add_shape(1, Inches(0), top, SLIDE_W, Pt(3))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()


def slide_title_block(slide, title, subtitle=None):
    add_textbox(slide, title,
                Inches(0.6), Inches(0.2), Inches(12), Inches(0.85),
                font_size=32, bold=True, colour=ACCENT, align=PP_ALIGN.LEFT)
    add_accent_bar(slide)
    if subtitle:
        add_textbox(slide, subtitle,
                    Inches(0.6), Inches(1.25), Inches(12), Inches(0.5),
                    font_size=16, colour=LIGHT_GREY, italic=True)


def add_ref(slide, text):
    add_textbox(slide, text,
                Inches(0.6), Inches(6.9), Inches(12), Inches(0.5),
                font_size=11, colour=LIGHT_GREY, italic=True)


# ── Build presentation ────────────────────────────────────────────────────────

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]   # completely blank layout


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title / Introduction
# ══════════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(blank)
set_bg(s1, DARK_BG)

add_textbox(s1, "GDPR for ICT Developers",
            Inches(1), Inches(1.8), Inches(11), Inches(1.2),
            font_size=44, bold=True, colour=ACCENT, align=PP_ALIGN.CENTER)

add_textbox(s1, "Key Principles, Obligations & Secure Implementation",
            Inches(1), Inches(3.0), Inches(11), Inches(0.7),
            font_size=22, colour=WHITE, align=PP_ALIGN.CENTER)

add_textbox(s1, "Case Study: GDPR-Compliant Student Registration API",
            Inches(1), Inches(3.75), Inches(11), Inches(0.6),
            font_size=18, colour=LIGHT_GREY, italic=True, align=PP_ALIGN.CENTER)

add_textbox(s1, "Chioma Debbie Okoye  |  MSc Cybersecurity, Privacy and Trust",
            Inches(1), Inches(5.2), Inches(11), Inches(0.5),
            font_size=15, colour=LIGHT_GREY, align=PP_ALIGN.CENTER)

add_textbox(s1, "Reference: Regulation (EU) 2016/679 — General Data Protection Regulation",
            Inches(1), Inches(6.8), Inches(11), Inches(0.4),
            font_size=11, colour=LIGHT_GREY, italic=True, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — What is GDPR & Why Developers Must Care
# ══════════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(blank)
set_bg(s2, DARK_BG)
slide_title_block(s2, "What is GDPR & Why Developers Must Care",
                  "Slide 2 of 6  ·  ~1 minute")

add_bullet_box(s2, [
    "GDPR (EU) 2016/679 — applies to any system processing personal data of EU residents",
    "Personal data: name, address, date of birth, phone number, health/disability information",
    "Disability data = Special Category (Article 9) — requires explicit consent",
    "Developers are legally accountable — 'privacy by design' is mandatory (Article 25)",
    "Fines up to €20 million or 4% of global annual turnover for non-compliance",
    "Our system collects all of the above — GDPR compliance is not optional",
], Inches(0.6), Inches(1.6), Inches(12), Inches(4.5), font_size=17)

add_ref(s2, "European Parliament (2016) Regulation (EU) 2016/679 (GDPR). Available at: https://gdpr-info.eu")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Data Minimisation & Lawful Basis
# ══════════════════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(blank)
set_bg(s3, DARK_BG)
slide_title_block(s3, "Data Minimisation & Lawful Basis",
                  "Slide 3 of 6  ·  ~1 minute")

add_bullet_box(s3, [
    "Article 5 — Data must be adequate, relevant and limited to what is necessary",
    "Only collect fields you actually use — our API collects 5 personal fields",
    "Article 6 — Processing must have a lawful basis (e.g. consent, contract)",
    "Article 9 — Disability data requires explicit, informed consent before collection",
    "Developers must implement consent mechanisms — not just store the data",
    "Data retention: define how long records are kept and delete when no longer needed",
], Inches(0.6), Inches(1.6), Inches(12), Inches(4.5), font_size=17)

add_ref(s3, "ICO (2021) Guide to the UK GDPR — Lawful Basis. Available at: https://ico.org.uk/for-organisations/guide-to-data-protection/")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Encryption at Rest: Protecting Personal Data
# ══════════════════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(blank)
set_bg(s4, DARK_BG)
slide_title_block(s4, "Encryption at Rest — Protecting Personal Data",
                  "Slide 4 of 6  ·  ~1 minute")

add_bullet_box(s4, [
    "Article 32 — Implement appropriate technical measures including encryption",
    "Our system uses AES-256-CBC to encrypt all personal data fields before storage",
    "Each field gets a unique random IV — same value encrypts differently every time",
    "Encryption key stored in OS keyring — never in the database or source code",
], Inches(0.6), Inches(1.6), Inches(7.8), Inches(3.0), font_size=16)

add_code_box(s4,
    "# api/crypto.py\n"
    "def encrypt_field(plaintext: str):\n"
    "    key = get_encryption_key()   # from OS keyring\n"
    "    iv  = os.urandom(16)         # unique per field\n"
    "    # AES-256-CBC encrypt → base64\n"
    "    return ciphertext_b64, iv_b64\n\n"
    "# MongoDB stores:\n"
    "# fullName:    'aGVsbG8gd29ybGQ=...'\n"
    "# fullName_iv: 'dGhpcyBpcyBhbiBJVg=='",
    Inches(8.5), Inches(1.6), Inches(4.6), Inches(3.5), font_size=12)

add_ref(s4, "Stallings, W. (2017) Cryptography and Network Security, 7th edn. Pearson.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Password Hashing & Token Security
# ══════════════════════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(blank)
set_bg(s5, DARK_BG)
slide_title_block(s5, "Password Hashing & Token Security",
                  "Slide 5 of 6  ·  ~1 minute")

add_bullet_box(s5, [
    "Passwords must NEVER be stored in plaintext — Article 32 requires appropriate security",
    "Our system uses bcrypt with work factor 12 — slow by design, resists brute force",
    "bcrypt embeds a random salt — same password hashes differently every time",
    "Session tokens hashed with SHA-256 before storage — token in DB ≠ token sent to user",
    "If DB is breached: attacker gets hashes, not passwords or valid tokens",
], Inches(0.6), Inches(1.6), Inches(7.8), Inches(3.2), font_size=16)

add_code_box(s5,
    "# Stored in MongoDB after registration:\n"
    "password: '$2b$12$eImiTXuWVxfM...'\n\n"
    "# Stored after login:\n"
    "token: 'a3f5c2d1e8b7...  (SHA-256 hash)'\n\n"
    "# run_hacker.py output — attacker sees:\n"
    "# password: $2b$12$... (bcrypt)\n"
    "# token:    9f86d081... (SHA-256)\n"
    "# fullName: aGVsbG8=   (ciphertext)",
    Inches(8.5), Inches(1.6), Inches(4.6), Inches(3.5), font_size=12)

add_ref(s5, "NIST (2017) SP 800-63B: Digital Identity Guidelines. Gaithersburg: NIST.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Breach Notification & Developer Responsibilities
# ══════════════════════════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(blank)
set_bg(s6, DARK_BG)
slide_title_block(s6, "Breach Notification & Developer Responsibilities",
                  "Slide 6 of 6  ·  ~1 minute")

add_bullet_box(s6, [
    "Article 33 — Notify supervisory authority within 72 hours of discovering a breach",
    "Article 34 — Notify affected individuals if breach poses high risk to their rights",
    "Encryption reduces risk: if data is encrypted and key is safe, breach may not require individual notification",
    "Developers must maintain records of processing activities (Article 30)",
    "Privacy by Design (Article 25): security must be built in from the start — not added later",
], Inches(0.6), Inches(1.6), Inches(12), Inches(3.6), font_size=17)

# Summary box
shape = s6.shapes.add_shape(1, Inches(0.6), Inches(5.3), Inches(12), Inches(1.3))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0x00, 0x3A, 0x5C)
shape.line.color.rgb = ACCENT
shape.line.width = Pt(1.5)
tf = shape.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "Key takeaway: Encryption + Hashing + Key Management = GDPR-compliant by design"
run.font.size = Pt(17)
run.font.bold = True
run.font.color.rgb = YELLOW

add_ref(s6,
    "References: GDPR (EU) 2016/679 | ICO (2021) ico.org.uk | NIST SP 800-63B | "
    "Stallings (2017) Cryptography & Network Security, Pearson")


# ── Save ──────────────────────────────────────────────────────────────────────
out = "GDPR_Screencast.pptx"
prs.save(out)
print(f"Saved: {out}")
